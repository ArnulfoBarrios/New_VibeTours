# generate-presentation.py
# ----------------------------------------------------
# VibeTours Presentation Builder Script
# Designed for high-fidelity dark-mode UI presentation.
# Output: VibeTours_Presentation.pptx and VibeTours_Presentation.pdf
# ----------------------------------------------------

import os
import sys
import math
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageOps

# Define color palette for PIL (Tuples)
PIL_BG = (9, 11, 16)         # #090B10
PIL_CARD = (28, 28, 30)      # #1C1C1E
PIL_BORDER = (56, 56, 58)    # #38383A
PIL_PRIMARY = (0, 122, 255)  # #007AFF (Blue)
PIL_ACCENT = (175, 82, 222)  # #AF52DE (Purple)
PIL_SUCCESS = (16, 185, 129) # #10B981 (Green)
PIL_TEXT_PRIMARY = (255, 255, 255) # #FFFFFF
PIL_TEXT_SECONDARY = (184, 184, 184) # #B8B8B8

# Directories
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCREENSHOTS_DIR = os.path.join(os.path.dirname(PROJECT_ROOT), "Capturas de pantalla")
ASSETS_DIR = os.path.join(PROJECT_ROOT, "assets", "images")
TEMP_DIR = os.path.join(PROJECT_ROOT, "temp_presentation_assets")
OUTPUT_PPTX = os.path.join(PROJECT_ROOT, "VibeTours_Presentation.pptx")
OUTPUT_PDF = os.path.join(PROJECT_ROOT, "VibeTours_Presentation.pdf")

# Ensure temp directory exists
os.makedirs(TEMP_DIR, exist_ok=True)

# Standard Windows Font Paths
WIN_FONTS = {
    "bold": [
        "C:\\Windows\\Fonts\\segoeuib.ttf", # Segoe UI Bold
        "C:\\Windows\\Fonts\\arialbd.ttf"  # Arial Bold
    ],
    "medium": [
        "C:\\Windows\\Fonts\\segoeuisl.ttf", # Segoe UI Semibold
        "C:\\Windows\\Fonts\\segoeui.ttf",   # Segoe UI Regular
        "C:\\Windows\\Fonts\\arial.ttf"      # Arial Regular
    ],
    "regular": [
        "C:\\Windows\\Fonts\\segoeui.ttf",   # Segoe UI Regular
        "C:\\Windows\\Fonts\\arial.ttf"      # Arial Regular
    ]
}

def get_font(font_type, size):
    """Load Segoe UI or Arial from Windows font folder or fall back to default."""
    for font_path in WIN_FONTS.get(font_type, []):
        if os.path.exists(font_path):
            try:
                return ImageFont.truetype(font_path, size)
            except Exception as e:
                print(f"Failed to load font {font_path}: {e}")
    try:
        return ImageFont.load_default()
    except:
        return None

def wrap_text(text, font, max_width):
    """Wrap text to fit within a maximum width in pixels."""
    lines = []
    words = text.split()
    current_line = []
    for word in words:
        test_line = ' '.join(current_line + [word])
        w = font.getlength(test_line)
        if w <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
            current_line = [word]
    if current_line:
        lines.append(' '.join(current_line))
    return lines

def get_screenshot_path(name):
    """Helper to locate screenshot by name in standard directory."""
    for folder in [SCREENSHOTS_DIR, os.path.join(PROJECT_ROOT, "Capturas de pantalla")]:
        p = os.path.join(folder, name)
        if os.path.exists(p):
            return p
    return None

# =========================================================================
# GRAPHICAL HELPERS
# =========================================================================

def draw_radial_glow(draw, center, radius, color, alpha_max):
    """Draw a soft radial gradient glow using math/alpha pixels directly on an image."""
    cx, cy = center
    overlay = Image.new("RGBA", (radius * 2, radius * 2), (0, 0, 0, 0))
    odraw = ImageDraw.Draw(overlay)
    
    for r in range(radius, 0, -4):
        ratio = 1.0 - (r / radius)
        alpha = int(alpha_max * (1.0 - ratio * ratio))
        odraw.ellipse([radius - r, radius - r, radius + r, radius + r], fill=color + (alpha,))
        
    return overlay

def create_iphone_mockup(screenshot_img, output_w=320, output_h=690):
    """Wrap a resized screenshot in a premium, detailed vector iPhone mockup with dynamic island."""
    canvas_w = output_w + 60
    canvas_h = output_h + 60
    
    img = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    
    cx = canvas_w // 2
    cy = canvas_h // 2
    
    # Draw dark shadow
    shadow_img = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
    sdraw = ImageDraw.Draw(shadow_img)
    sdraw.rounded_rectangle(
        [cx - output_w//2 + 8, cy - output_h//2 + 10, cx + output_w//2 + 8, cy + output_h//2 + 10],
        radius=40, fill=(0, 0, 0, 160)
    )
    shadow_img = shadow_img.filter(ImageFilter.GaussianBlur(15))
    img.paste(shadow_img, (0, 0), shadow_img)
    
    # Frame dimensions
    fx1 = cx - output_w // 2
    fy1 = cy - output_h // 2
    fx2 = cx + output_w // 2
    fy2 = cy + output_h // 2
    
    # Screen dimensions
    screen_padding = 8
    sx1 = fx1 + screen_padding
    sy1 = fy1 + screen_padding
    sx2 = fx2 - screen_padding
    sy2 = fy2 - screen_padding
    sw = sx2 - sx1
    sh = sy2 - sy1
    
    # Resize and crop the screenshot
    ss_resized = screenshot_img.resize((sw, sh), Image.Resampling.LANCZOS)
    
    # Rounded corner mask
    mask = Image.new("L", (sw, sh), 0)
    mdraw = ImageDraw.Draw(mask)
    mdraw.rounded_rectangle([0, 0, sw, sh], radius=32, fill=255)
    
    screen_canvas = Image.new("RGBA", (sw, sh), (0, 0, 0, 255))
    screen_canvas.paste(ss_resized, (0, 0), mask)
    
    # Draw Dynamic Island
    sdraw_ctx = ImageDraw.Draw(screen_canvas)
    di_w = 100
    di_h = 22
    di_x1 = sw // 2 - di_w // 2
    di_y1 = 18
    sdraw_ctx.rounded_rectangle([di_x1, di_y1, di_x1 + di_w, di_y1 + di_h], radius=11, fill=(0, 0, 0, 255))
    
    img.paste(screen_canvas, (sx1, sy1), screen_canvas)
    
    # Draw outer metallic iPhone frame
    draw.rounded_rectangle([fx1, fy1, fx2, fy2], radius=40, outline=PIL_BORDER, width=screen_padding)
    draw.rounded_rectangle([fx1 - 1, fy1 - 1, fx2 + 1, fy2 + 1], radius=41, outline=(255, 255, 255, 30), width=1)
    
    return img

def generate_report_screenshot():
    """Generate a custom screen representing a tour with the iOS-style report bottom sheet overlaid."""
    base_path = get_screenshot_path("Pantallazo8.jpeg")
    if not base_path:
        base_img = Image.new("RGBA", (540, 1200), PIL_BG)
    else:
        base_img = Image.open(base_path).convert("RGBA")
        
    base_img = base_img.resize((540, 1200), Image.Resampling.LANCZOS)
    draw = ImageDraw.Draw(base_img)
    
    # Add dark screen overlay
    overlay = Image.new("RGBA", (540, 1200), (0, 0, 0, 140))
    base_img.paste(overlay, (0, 0), overlay)
    
    # Draw bottom sheet card
    card_y = 680
    card_h = 520
    draw.rounded_rectangle([15, card_y, 525, card_y + card_h], radius=36, fill=PIL_CARD, outline=PIL_BORDER, width=3)
    
    # Handle bar
    draw.rounded_rectangle([230, card_y + 16, 310, card_y + 22], radius=3, fill=(90, 90, 92, 255))
    
    font_bold = get_font("bold", 36)
    font_regular = get_font("regular", 26)
    draw.text((45, card_y + 50), "Reportar Tour", font=font_bold, fill=PIL_TEXT_PRIMARY)
    
    # Options
    options = ["Spam / Publicidad molesta", "Contenido inapropiado u ofensivo", "Falta de seguridad en el recorrido", "Información falsa o errónea", "Otro motivo"]
    curr_y = card_y + 120
    for idx, opt in enumerate(options):
        is_selected = (idx == 1)
        rx = 45
        ry = curr_y + 8
        if is_selected:
            draw.ellipse([rx, ry, rx + 30, ry + 30], fill=PIL_PRIMARY, outline=PIL_PRIMARY, width=2)
            draw.ellipse([rx + 8, ry + 8, rx + 22, ry + 22], fill=PIL_TEXT_PRIMARY)
        else:
            draw.ellipse([rx, ry, rx + 30, ry + 30], outline=PIL_TEXT_SECONDARY, width=2)
            
        draw.text((95, curr_y + 5), opt, font=font_regular, fill=PIL_TEXT_PRIMARY if is_selected else PIL_TEXT_SECONDARY)
        curr_y += 60
        
    # Draw button
    btn_y = card_y + 420
    draw.rounded_rectangle([45, btn_y, 495, btn_y + 65], radius=16, fill=PIL_PRIMARY)
    
    btn_text = "Enviar Reporte"
    btn_w = font_bold.getlength(btn_text)
    draw.text((270 - btn_w // 2, btn_y + 15), btn_text, font=font_bold, fill=PIL_TEXT_PRIMARY)
    
    return base_img

def generate_pqrs_screenshot():
    """Programmatically draw a premium representation of the Support (PQRS) screen matching help_center_screen.dart."""
    img = Image.new("RGBA", (540, 1200), PIL_BG)
    draw = ImageDraw.Draw(img)
    
    glow = draw_radial_glow(draw, (270, 700), 450, PIL_PRIMARY, 25)
    img.paste(glow, (270 - 450, 700 - 450), glow)
    
    font_bold = get_font("bold", 38)
    font_medium = get_font("medium", 28)
    font_regular = get_font("regular", 22)
    font_small = get_font("regular", 18)
    
    # Top header bar
    draw.text((36, 54), "VibeTours", font=font_bold, fill=PIL_TEXT_PRIMARY)
    
    # Support Pill
    draw.rounded_rectangle([250, 52, 380, 94], radius=21, fill=(0, 122, 255, 38), outline=PIL_PRIMARY, width=2)
    draw.text((276, 60), "Soporte", font=font_medium, fill=PIL_PRIMARY)
    
    # Close icon (X)
    draw.ellipse([460, 52, 502, 94], fill=PIL_CARD, outline=PIL_BORDER, width=2)
    draw.line([473, 65, 489, 81], fill=PIL_TEXT_PRIMARY, width=3)
    draw.line([489, 65, 473, 81], fill=PIL_TEXT_PRIMARY, width=3)
    
    # Tabs
    draw.text((36, 140), "Crear Solicitud", font=font_medium, fill=PIL_PRIMARY)
    draw.rectangle([36, 180, 270, 185], fill=PIL_PRIMARY)
    draw.text((310, 140), "Historial", font=font_medium, fill=PIL_TEXT_SECONDARY)
    
    # Main Form Card
    form_y = 210
    form_h = 680
    draw.rounded_rectangle([30, form_y, 510, form_y + form_h], radius=34, fill=PIL_CARD, outline=PIL_BORDER, width=2)
    
    # Title & description inside card
    draw.text((60, form_y + 40), "Centro de Ayuda", font=font_bold, fill=PIL_TEXT_PRIMARY)
    desc = "Envía Peticiones, Quejas, Reclamos y Sugerencias de soporte para su procesamiento."
    wrapped_desc = wrap_text(desc, font_regular, 400)
    dy = form_y + 90
    for line in wrapped_desc:
        draw.text((60, dy), line, font=font_regular, fill=PIL_TEXT_SECONDARY)
        dy += 32
        
    # Input 1: Tipo de Solicitud
    label_y = form_y + 190
    draw.text((60, label_y), "Tipo de solicitud", font=font_medium, fill=PIL_TEXT_PRIMARY)
    box_y = label_y + 40
    draw.rounded_rectangle([60, box_y, 480, box_y + 70], radius=16, fill=PIL_BG, outline=PIL_BORDER, width=2)
    draw.text((80, box_y + 20), "Sugerencia", font=font_regular, fill=PIL_TEXT_PRIMARY)
    draw.polygon([(430, box_y + 30), (450, box_y + 30), (440, box_y + 45)], fill=PIL_TEXT_SECONDARY)
    
    # Input 2: Asunto
    label2_y = box_y + 100
    draw.text((60, label2_y), "Asunto", font=font_medium, fill=PIL_TEXT_PRIMARY)
    box2_y = label2_y + 40
    draw.rounded_rectangle([60, box2_y, 480, box2_y + 70], radius=16, fill=PIL_BG, outline=PIL_BORDER, width=2)
    draw.text((80, box2_y + 20), "Sugerencia sobre la interfaz de mapa", font=font_regular, fill=PIL_TEXT_PRIMARY)
    
    # Input 3: Mensaje
    label3_y = box2_y + 100
    draw.text((60, label3_y), "Mensaje", font=font_medium, fill=PIL_TEXT_PRIMARY)
    box3_y = label3_y + 40
    draw.rounded_rectangle([60, box3_y, 480, box3_y + 150], radius=16, fill=PIL_BG, outline=PIL_BORDER, width=2)
    draw.text((80, box3_y + 20), "Me gustaría proponer una mejora en la...", font=font_regular, fill=PIL_TEXT_SECONDARY)
    
    # Submit Button
    btn_y = box3_y + 180
    draw.rounded_rectangle([60, btn_y, 480, btn_y + 70], radius=16, fill=PIL_PRIMARY)
    btn_text = "Enviar PQRS"
    btn_w = font_bold.getlength(btn_text)
    draw.text((270 - btn_w // 2, btn_y + 15), btn_text, font=font_bold, fill=PIL_TEXT_PRIMARY)
    
    # Bottom Info Cards
    card2_y = form_y + form_h + 30
    draw.rounded_rectangle([30, card2_y, 255, card2_y + 180], radius=24, fill=PIL_CARD, outline=PIL_BORDER, width=2)
    draw.ellipse([50, card2_y + 20, 90, card2_y + 60], outline=PIL_PRIMARY, width=3)
    draw.line([70, card2_y + 30, 70, card2_y + 40], fill=PIL_PRIMARY, width=3)
    draw.line([70, card2_y + 40, 80, card2_y + 40], fill=PIL_PRIMARY, width=3)
    draw.text((50, card2_y + 75), "Respuesta rápida", font=font_small, fill=PIL_TEXT_PRIMARY)
    draw.text((50, card2_y + 105), "En menos de 24h", font=font_small, fill=PIL_TEXT_SECONDARY)
    
    draw.rounded_rectangle([285, card2_y, 510, card2_y + 180], radius=24, fill=PIL_CARD, outline=PIL_BORDER, width=2)
    draw.polygon([(325, card2_y + 20), (345, card2_y + 15), (365, card2_y + 20), (365, card2_y + 40), (345, card2_y + 60), (325, card2_y + 40)], outline=PIL_PRIMARY, width=3)
    draw.text((305, card2_y + 75), "Conexión segura", font=font_small, fill=PIL_TEXT_PRIMARY)
    draw.text((305, card2_y + 105), "Supabase SSL", font=font_small, fill=PIL_TEXT_SECONDARY)
    
    return img

# =========================================================================
# DRAWING VECTOR BRAND LOGOS LOCALLY
# =========================================================================

def draw_vector_logo_offline(name, w=96, h=96):
    """Generate high-fidelity stylized technology logos offline using Pillow vector primitives."""
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    cx, cy = w // 2, h // 2
    
    if name == "flutter":
        # Draw Flutter chevrons
        # Points of top/right chevron
        draw.polygon([(cx - 10, cy - 35), (cx + 25, cy), (cx + 10, cy + 15), (cx - 25, cy - 20)], fill=(69, 209, 253, 255))
        # Points of bottom chevron
        draw.polygon([(cx - 25, cy - 20), (cx + 10, cy + 15), (cx - 10, cy + 35), (cx - 28, cy + 17)], fill=(2, 86, 155, 255))
        # Center triangle shading
        draw.polygon([(cx - 10, cy), (cx + 10, cy + 15), (cx - 25, cy - 20)], fill=(3, 169, 244, 255))
        
    elif name == "supabase":
        # Lightning Bolt logo
        # Lightning shape
        pts = [
            (cx - 5, cy - 35), (cx + 20, cy - 5), (cx + 2, cy - 5),
            (cx + 10, cy + 35), (cx - 15, cy + 5), (cx - 2, cy + 5)
        ]
        draw.polygon(pts, fill=(62, 207, 142, 255))
        
    elif name == "firebase":
        # Flame/Shield shape
        draw.polygon([(cx, cy - 38), (cx + 28, cy + 10), (cx + 18, cy + 30), (cx - 18, cy + 30), (cx - 28, cy + 10)], fill=(255, 202, 40, 255))
        draw.polygon([(cx, cy - 38), (cx + 22, cy + 15), (cx, cy + 30), (cx - 22, cy + 15)], fill=(245, 124, 0, 255))
        draw.polygon([(cx, cy - 38), (cx, cy + 30), (cx + 12, cy + 22)], fill=(226, 135, 67, 255))
        
    elif name == "nodejs":
        # Green Hexagon
        side = 28
        pts = []
        for i in range(6):
            angle = math.radians(i * 60)
            pts.append((cx + side * math.sin(angle), cy + side * math.cos(angle)))
        draw.polygon(pts, outline=(51, 153, 51, 255), fill=(51, 153, 51, 60), width=4)
        # Stylized Node text/structure
        draw.ellipse([cx - 8, cy - 8, cx + 8, cy + 8], fill=(255, 255, 255, 200))
        
    elif name == "openai":
        # OpenAI spiral logo
        # Draw central circle and surrounding leaves using rotated ellipses
        for i in range(6):
            angle = i * 60
            # Draw a capsule rotated around center
            capsule = Image.new("RGBA", (w, h), (0, 0, 0, 0))
            cdraw = ImageDraw.Draw(capsule)
            cdraw.rounded_rectangle([cx - 6, cy - 32, cx + 6, cy - 2], radius=6, fill=PIL_TEXT_PRIMARY)
            capsule = capsule.rotate(angle, center=(cx, cy))
            img.paste(capsule, (0, 0), capsule)
            
    elif name == "tomtom":
        # Red abstract logo
        draw.ellipse([cx - 28, cy - 28, cx + 28, cy + 28], fill=(223, 27, 18, 255))
        # Draw white stylized center
        draw.ellipse([cx - 10, cy - 10, cx + 10, cy + 10], fill=PIL_TEXT_PRIMARY)
        draw.line([cx - 18, cy, cx + 18, cy], fill=PIL_TEXT_PRIMARY, width=4)
        
    elif name == "maplibre":
        # Map shape
        # Draw three columns (folded map)
        draw.polygon([(cx - 24, cy - 22), (cx - 8, cy - 32), (cx - 8, cy + 18), (cx - 24, cy + 28)], fill=(44, 62, 80, 255), outline=PIL_BORDER)
        draw.polygon([(cx - 8, cy - 32), (cx + 8, cy - 22), (cx + 8, cy + 28), (cx - 8, cy + 18)], fill=(62, 207, 142, 255), outline=PIL_BORDER)
        draw.polygon([(cx + 8, cy - 22), (cx + 24, cy - 32), (cx + 24, cy + 18), (cx + 8, cy + 28)], fill=(44, 62, 80, 255), outline=PIL_BORDER)
        
    return img

def draw_vector_icon(name, size=64):
    """Draw clean outline vector icons programmatically to represent stations."""
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    
    mid = size // 2
    pad = 8
    
    if name == "login":
        draw.ellipse([pad + 4, pad + 4, mid + 2, mid + 2], outline=PIL_TEXT_PRIMARY, width=4)
        draw.line([mid - 2, mid - 2, size - pad, size - pad], fill=PIL_TEXT_PRIMARY, width=4)
        draw.line([size - pad - 6, size - pad, size - pad, size - pad - 6], fill=PIL_TEXT_PRIMARY, width=4)
        draw.line([size - pad - 12, size - pad - 6, size - pad - 6, size - pad - 12], fill=PIL_TEXT_PRIMARY, width=4)
    elif name == "create":
        draw.polygon([
            (mid, pad), (mid + 5, mid - 5), (size - pad, mid), (mid + 5, mid + 5),
            (mid, size - pad), (mid - 5, mid + 5), (pad, mid), (mid - 5, mid - 5)
        ], outline=PIL_TEXT_PRIMARY, width=4)
    elif name == "publish":
        draw.ellipse([pad, pad, size - pad, size - pad], outline=PIL_TEXT_PRIMARY, width=4)
        draw.ellipse([pad + 12, pad, size - pad - 12, size - pad], outline=PIL_TEXT_PRIMARY, width=2)
        draw.line([pad, mid, size - pad, mid], fill=PIL_TEXT_PRIMARY, width=2)
        draw.line([mid, pad, mid, size - pad], fill=PIL_TEXT_PRIMARY, width=2)
    elif name == "report":
        draw.polygon([(mid, pad), (pad, size - pad), (size - pad, size - pad)], outline=PIL_TEXT_PRIMARY, width=4)
        draw.line([mid, mid - 10, mid, mid + 6], fill=PIL_TEXT_PRIMARY, width=4)
        draw.ellipse([mid - 2, mid + 12, mid + 2, mid + 16], fill=PIL_TEXT_PRIMARY)
    elif name == "support":
        draw.rounded_rectangle([pad, pad + 6, size - pad, size - pad - 6], radius=8, outline=PIL_TEXT_PRIMARY, width=4)
        draw.line([pad + 6, pad + 12, mid, mid + 2], fill=PIL_TEXT_PRIMARY, width=4)
        draw.line([mid, mid + 2, size - pad - 6, pad + 12], fill=PIL_TEXT_PRIMARY, width=4)
    elif name == "delete":
        draw.line([pad, pad + 10, size - pad, pad + 10], fill=PIL_TEXT_PRIMARY, width=4)
        draw.rounded_rectangle([pad + 8, pad + 10, size - pad - 8, size - pad], radius=6, outline=PIL_TEXT_PRIMARY, width=4)
        draw.line([mid - 6, pad + 18, mid - 6, size - pad - 8], fill=PIL_TEXT_PRIMARY, width=3)
        draw.line([mid + 6, pad + 18, mid + 6, size - pad - 8], fill=PIL_TEXT_PRIMARY, width=3)
        draw.line([mid - 12, pad + 10, mid - 12, pad + 4], fill=PIL_TEXT_PRIMARY, width=3)
        draw.line([mid - 12, pad + 4, mid + 12, pad + 4], fill=PIL_TEXT_PRIMARY, width=3)
        draw.line([mid + 12, pad + 4, mid + 12, pad + 10], fill=PIL_TEXT_PRIMARY, width=3)
        
    return img

# =========================================================================
# ASSET COMPOSITION & PROCESSING
# =========================================================================

def process_all_screenshots():
    """Load screenshots, generate report/pqrs screens, and wrap all in iPhone mockups."""
    print("Processing screenshots and wrapping in premium iPhone mockups...")
    
    # Generate vector logos locally offline
    for tech in ["flutter", "supabase", "firebase", "nodejs", "openai", "tomtom", "maplibre"]:
        dest = os.path.join(TEMP_DIR, f"logo_{tech}.png")
        logo_img = draw_vector_logo_offline(tech, 96, 96)
        logo_img.save(dest)
        
    mapping = {
        "1_login": "Pantallazo1.jpeg",
        "2_chat_ia": "Pantallazo3.jpeg",
        "3_publicar": "Pantallazo7.jpeg",
        "6_perfil": "Pantallazo5.jpeg",
    }
    
    # Process standard screenshots
    for key, ss_name in mapping.items():
        p = get_screenshot_path(ss_name)
        if p:
            ss_img = Image.open(p).convert("RGBA")
            mockup = create_iphone_mockup(ss_img, output_w=320, output_h=690)
            mockup.save(os.path.join(TEMP_DIR, f"mockup_{key}.png"))
        else:
            print(f"Warning: Screenshot {ss_name} not found, generating a placeholder.")
            placeholder = Image.new("RGBA", (540, 1200), PIL_CARD)
            pdraw = ImageDraw.Draw(placeholder)
            pdraw.text((50, 400), f"VibeTours\n{key}", font=get_font("bold", 38), fill=PIL_PRIMARY)
            mockup = create_iphone_mockup(placeholder, output_w=320, output_h=690)
            mockup.save(os.path.join(TEMP_DIR, f"mockup_{key}.png"))
            
    # Process custom screen for Station 4 (Reportar)
    report_ss = generate_report_screenshot()
    mockup_report = create_iphone_mockup(report_ss, output_w=320, output_h=690)
    mockup_report.save(os.path.join(TEMP_DIR, "mockup_4_reportar.png"))
    
    # Process custom screen for Station 5 (PQRS)
    pqrs_ss = generate_pqrs_screenshot()
    mockup_pqrs = create_iphone_mockup(pqrs_ss, output_w=320, output_h=690)
    mockup_pqrs.save(os.path.join(TEMP_DIR, "mockup_5_pqrs.png"))
    
    # Save outline icons
    icon_names = ["login", "create", "publish", "report", "support", "delete"]
    for ic in icon_names:
        img_ic = draw_vector_icon(ic, size=64)
        img_ic.save(os.path.join(TEMP_DIR, f"icon_{ic}.png"))

# =========================================================================
# SLIDE BG GENERATORS (PILLOW)
# =========================================================================

def draw_bezier_segment(draw, p0, p1, p2, p3, steps=100):
    pts = []
    for i in range(steps + 1):
        t = i / steps
        x = (1-t)**3 * p0[0] + 3*(1-t)**2 * t * p1[0] + 3*(1-t) * t**2 * p2[0] + t**3 * p3[0]
        y = (1-t)**3 * p0[1] + 3*(1-t)**2 * t * p1[1] + 3*(1-t) * t**2 * p2[1] + t**3 * p3[1]
        pts.append((x, y))
    return pts

def generate_cover_background():
    img = Image.new("RGBA", (1920, 1080), PIL_BG)
    draw = ImageDraw.Draw(img)
    
    glow_blue = draw_radial_glow(draw, (1400, 300), 700, PIL_PRIMARY, 40)
    img.paste(glow_blue, (1400 - 700, 300 - 700), glow_blue)
    
    glow_purple = draw_radial_glow(draw, (400, 800), 700, PIL_ACCENT, 45)
    img.paste(glow_purple, (400 - 700, 800 - 700), glow_purple)
    
    # Grid lines
    for y in range(100, 1080, 120):
        draw.line([0, y, 1920, y], fill=(255, 255, 255, 10), width=1)
    for x in range(100, 1920, 150):
        draw.line([x, 0, x, 1080], fill=(255, 255, 255, 10), width=1)
        
    # Particles
    import random
    random.seed(42)
    for _ in range(70):
        px = random.randint(50, 1870)
        py = random.randint(50, 1030)
        size = random.randint(2, 6)
        alpha = random.randint(30, 160)
        draw.ellipse([px - size*2, py - size*2, px + size*2, py + size*2], fill=(255, 255, 255, alpha // 3))
        draw.ellipse([px - size//2, py - size//2, px + size//2, py + size//2], fill=(255, 255, 255, alpha))
        
    img.save(os.path.join(TEMP_DIR, "bg_cover.png"))

def generate_roadmap_background():
    img = Image.new("RGBA", (1920, 1080), PIL_BG)
    draw = ImageDraw.Draw(img)
    
    glow_blue = draw_radial_glow(draw, (960, 540), 900, PIL_PRIMARY, 20)
    img.paste(glow_blue, (960 - 900, 540 - 900), glow_blue)
    
    glow_purple = draw_radial_glow(draw, (1600, 800), 600, PIL_ACCENT, 25)
    img.paste(glow_purple, (1600 - 600, 800 - 600), glow_purple)
    
    stations_coords = [
        (180, 520),
        (480, 720),
        (780, 520),
        (1080, 720),
        (1380, 520),
        (1680, 720)
    ]
    
    path_points = []
    for i in range(len(stations_coords) - 1):
        p0 = stations_coords[i]
        p3 = stations_coords[i+1]
        p1 = (p0[0] + 150, p0[1])
        p2 = (p3[0] - 150, p3[1])
        segment = draw_bezier_segment(draw, p0, p1, p2, p3, steps=60)
        path_points.extend(segment[:-1])
    path_points.append(stations_coords[-1])
    
    # Path glows
    for w, opac in [(24, 8), (14, 25), (6, 75)]:
        for k in range(len(path_points) - 1):
            draw.line([path_points[k], path_points[k+1]], fill=PIL_PRIMARY + (opac,), width=w)
            
    # Sharp core
    for k in range(len(path_points) - 1):
        draw.line([path_points[k], path_points[k+1]], fill=(255, 255, 255, 220), width=2)
        
    # Station nodes
    for idx, (sx, sy) in enumerate(stations_coords):
        color = PIL_PRIMARY if idx % 2 == 0 else PIL_ACCENT
        draw.ellipse([sx - 35, sy - 35, sx + 35, sy + 35], fill=color + (30,), outline=color + (90,), width=2)
        draw.ellipse([sx - 12, sy - 12, sx + 12, sy + 12], fill=(255, 255, 255, 255), outline=color, width=4)
        
        num_str = str(idx + 1)
        font_num = get_font("bold", 15)
        draw.text((sx - 5, sy - 10), num_str, font=font_num, fill=PIL_BG)
        
    img.save(os.path.join(TEMP_DIR, "bg_roadmap.png"))

def generate_dashboard_background():
    img = Image.new("RGBA", (1920, 1080), PIL_BG)
    draw = ImageDraw.Draw(img)
    
    glow_purple = draw_radial_glow(draw, (300, 200), 600, PIL_ACCENT, 30)
    img.paste(glow_purple, (300 - 600, 200 - 600), glow_purple)
    
    glow_blue = draw_radial_glow(draw, (1600, 800), 700, PIL_PRIMARY, 25)
    img.paste(glow_blue, (1600 - 700, 800 - 700), glow_blue)
    
    img.save(os.path.join(TEMP_DIR, "bg_dashboard.png"))

# =========================================================================
# PDF HIGH-RES RENDERER (PILLOW)
# =========================================================================

def draw_card_glass(draw, x1, y1, x2, y2, radius=24, border_w=2):
    draw.rounded_rectangle([x1, y1, x2, y2], radius=radius, fill=PIL_CARD + (240,), outline=PIL_BORDER + (255,), width=border_w)

def render_slide_0_image():
    img = Image.open(os.path.join(TEMP_DIR, "bg_cover.png")).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    font_bold_huge = get_font("bold", 96)
    font_med_mid = get_font("medium", 32)
    
    logo_path = os.path.join(ASSETS_DIR, "logo_light.png")
    if os.path.exists(logo_path):
        logo = Image.open(logo_path).convert("RGBA")
        logo = logo.resize((180, 180), Image.Resampling.LANCZOS)
        img.paste(logo, (960 - 90, 280), logo)
        
    title = "VibeTours."
    tw = font_bold_huge.getlength(title)
    draw.text((960 - tw // 2 + 3, 500 + 3), title, font=font_bold_huge, fill=(0, 0, 0, 120))
    draw.text((960 - tw // 2, 500), title, font=font_bold_huge, fill=PIL_TEXT_PRIMARY)
    
    slogan = "Planifica, crea y comparte experiencias impulsadas por Inteligencia Artificial."
    sw = font_med_mid.getlength(slogan)
    draw.text((960 - sw // 2, 630), slogan, font=font_med_mid, fill=PIL_TEXT_SECONDARY)
    
    draw.rounded_rectangle([960 - 150, 710, 960 + 150, 760], radius=25, fill=(0, 122, 255, 30), outline=PIL_PRIMARY, width=2)
    badge_text = "Apple Intelligence Standard"
    font_badge = get_font("bold", 18)
    bw = font_badge.getlength(badge_text)
    draw.text((960 - bw // 2, 726), badge_text, font=font_badge, fill=PIL_TEXT_PRIMARY)
    
    img.save(os.path.join(TEMP_DIR, "render_slide_0.png"))

def render_slide_1_image():
    img = Image.open(os.path.join(TEMP_DIR, "bg_roadmap.png")).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    font_title = get_font("bold", 54)
    font_subtitle = get_font("medium", 24)
    draw.text((80, 60), "El recorrido del usuario", font=font_title, fill=PIL_TEXT_PRIMARY)
    draw.text((80, 130), "CÓMO FUNCIONA VIBETOURS", font=font_subtitle, fill=PIL_PRIMARY)
    
    mockups = {
        1: ("mockup_1_login.png", 180, 180),
        2: ("mockup_2_chat_ia.png", 480, 810),
        3: ("mockup_3_publicar.png", 780, 180),
        4: ("mockup_4_reportar.png", 1080, 810),
        5: ("mockup_5_pqrs.png", 1380, 180),
        6: ("mockup_6_perfil.png", 1680, 810),
    }
    
    mw_canvas = 220
    mh_canvas = 400
    
    icons = {
        1: "icon_login.png",
        2: "icon_create.png",
        3: "icon_publish.png",
        4: "icon_report.png",
        5: "icon_support.png",
        6: "icon_delete.png"
    }
    
    stations_info = [
        {
            "num": 1,
            "title": "Ingreso Seguro",
            "desc": "Inicia sesión con Google, Correo o accede directo mediante el Modo Demo sin registros.",
            "pos": "top"
        },
        {
            "num": 2,
            "title": "Crear Tour (Manual / IA)",
            "desc": "Planifica manual o usa el asistente de IA dictando tus preferencias con micrófono por voz.",
            "pos": "bottom"
        },
        {
            "num": 3,
            "title": "Publicar en Catálogo",
            "desc": "Los tours IA se guardan como borradores temporales y pasan a moderación antes de publicarse.",
            "pos": "top"
        },
        {
            "num": 4,
            "title": "Reportar Contenido",
            "desc": "Filtro de seguridad. Reporta spam, contenido inapropiado o rutas peligrosas para moderación.",
            "pos": "bottom"
        },
        {
            "num": 5,
            "title": "Centro PQRS",
            "desc": "Soporte integrado en Supabase. Envía quejas o peticiones con respuesta en menos de 24h.",
            "pos": "top"
        },
        {
            "num": 6,
            "title": "Eliminación de Cuenta",
            "desc": "Garantía de privacidad. Controla tu cuenta, cierra sesión o elimina tu perfil permanentemente.",
            "pos": "bottom"
        }
    ]
    
    font_st_title = get_font("bold", 18)
    font_st_desc = get_font("regular", 13)
    
    for info in stations_info:
        idx = info["num"]
        mx, my_node = {
            1: (180, 440),
            2: (480, 640),
            3: (780, 440),
            4: (1080, 640),
            5: (1380, 440),
            6: (1680, 640)
        }[idx]
        
        file_name, sx, sy = mockups[idx]
        mimg = Image.open(os.path.join(TEMP_DIR, file_name))
        mimg_resized = mimg.resize((mw_canvas, mh_canvas), Image.Resampling.LANCZOS)
        px_mockup = sx - mw_canvas // 2
        py_mockup = sy - 40 if info["pos"] == "top" else sy - mh_canvas + 40
        img.paste(mimg_resized, (px_mockup, py_mockup), mimg_resized)
        
        tx_w = 300
        tx_h = 180
        px_card = sx - tx_w // 2
        py_card = 620 if info["pos"] == "top" else 200
        
        draw_card_glass(draw, px_card, py_card, px_card + tx_w, py_card + tx_h, radius=18)
        
        icon_path = os.path.join(TEMP_DIR, icons[idx])
        if os.path.exists(icon_path):
            ico = Image.open(icon_path).convert("RGBA")
            ico_resized = ico.resize((26, 26), Image.Resampling.LANCZOS)
            color_theme = PIL_PRIMARY if idx % 2 == 1 else PIL_ACCENT
            r, g, b, a = ico_resized.split()
            colored_ico = Image.merge("RGBA", (
                ImageOps.colorize(ImageOps.grayscale(ico_resized), (0,0,0), color_theme).split()[0],
                ImageOps.colorize(ImageOps.grayscale(ico_resized), (0,0,0), color_theme).split()[1],
                ImageOps.colorize(ImageOps.grayscale(ico_resized), (0,0,0), color_theme).split()[2],
                a
            ))
            img.paste(colored_ico, (px_card + 16, py_card + 14), colored_ico)
            
        draw.text((px_card + 52, py_card + 16), info["title"], font=font_st_title, fill=PIL_TEXT_PRIMARY)
        
        desc_lines = wrap_text(info["desc"], font_st_desc, tx_w - 32)
        dy = py_card + 54
        for line in desc_lines:
            draw.text((px_card + 16, dy), line, font=font_st_desc, fill=PIL_TEXT_SECONDARY)
            dy += 20
            
    img.save(os.path.join(TEMP_DIR, "render_slide_1.png"))

def render_slide_2_image():
    img = Image.open(os.path.join(TEMP_DIR, "bg_dashboard.png")).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    font_title = get_font("bold", 54)
    font_subtitle = get_font("medium", 24)
    draw.text((80, 60), "Estado actual del sistema", font=font_title, fill=PIL_TEXT_PRIMARY)
    draw.text((80, 130), "ESTRUCTURA DE LÍMITES Y CAPA DE INTEGRACIÓN TÉCNICA", font=font_subtitle, fill=PIL_PRIMARY)
    
    font_card_title = get_font("bold", 24)
    font_card_desc = get_font("regular", 16)
    
    # Card 1: Límite IA (Top Left)
    draw_card_glass(draw, 80, 200, 600, 540)
    draw.rounded_rectangle([110, 240, 116, 280], radius=3, fill=PIL_ACCENT)
    draw.text((130, 244), "Límite Demo IA", font=font_card_title, fill=PIL_TEXT_PRIMARY)
    desc1 = "Los usuarios invitados sin cuenta tienen un límite estricto de 2 generaciones de tours con IA. Al agotar el cupo de demostración local, el sistema solicita registrarse o iniciar sesión para continuar usando el planificador inteligente."
    lines1 = wrap_text(desc1, font_card_desc, 460)
    dy = 305
    for line in lines1:
        draw.text((110, dy), line, font=font_card_desc, fill=PIL_TEXT_SECONDARY)
        dy += 26
        
    # Card 2: Dependencias (Bottom Left)
    draw_card_glass(draw, 80, 570, 600, 910)
    draw.rounded_rectangle([110, 610, 116, 650], radius=3, fill=PIL_PRIMARY)
    draw.text((130, 614), "Dependencias de Mapas", font=font_card_title, fill=PIL_TEXT_PRIMARY)
    desc2 = "La cartografía y geocodificación se apoya en proveedores de fuentes abiertas: OpenStreetMap (OSM) para mapas, Photon y Nominatim para búsquedas de lugares, y Overpass API. Posee un mecanismo de contingencia para fallback automático."
    lines2 = wrap_text(desc2, font_card_desc, 460)
    dy = 675
    for line in lines2:
        draw.text((110, dy), line, font=font_card_desc, fill=PIL_TEXT_SECONDARY)
        dy += 26
        
    # Card 3: Moderación (Top Center)
    draw_card_glass(draw, 640, 200, 1160, 540)
    draw.rounded_rectangle([670, 240, 676, 280], radius=3, fill=PIL_PRIMARY)
    draw.text((690, 244), "Moderación de Tours", font=font_card_title, fill=PIL_TEXT_PRIMARY)
    desc3 = "Los tours creados mediante IA quedan guardados como borradores privados. El administrador debe revisarlos y aprobarlos desde su panel restringido antes de que aparezcan en el catálogo público general, previniendo spam."
    lines3 = wrap_text(desc3, font_card_desc, 460)
    dy = 305
    for line in lines3:
        draw.text((670, dy), line, font=font_card_desc, fill=PIL_TEXT_SECONDARY)
        dy += 26
        
    # Card 4: Administrador (Bottom Center)
    draw_card_glass(draw, 640, 570, 1160, 910)
    draw.rounded_rectangle([670, 610, 676, 650], radius=3, fill=PIL_ACCENT)
    draw.text((690, 614), "Acceso Administrador", font=font_card_title, fill=PIL_TEXT_PRIMARY)
    desc4 = "El panel de administración cuenta con seguridad granular. El acceso está restringido a un único ID de usuario o correo configurado estrictamente en las variables de entorno de la base de datos de Supabase, bloqueando intrusos."
    lines4 = wrap_text(desc4, font_card_desc, 460)
    dy = 675
    for line in lines4:
        draw.text((670, dy), line, font=font_card_desc, fill=PIL_TEXT_SECONDARY)
        dy += 26
        
    # Card 5: Tecnologías (Right - Column Complete)
    draw_card_glass(draw, 1200, 200, 1840, 910)
    draw.rounded_rectangle([1230, 240, 1236, 280], radius=3, fill=PIL_SUCCESS)
    draw.text((1250, 244), "Tecnologías & Servicios Integrados", font=font_card_title, fill=PIL_TEXT_PRIMARY)
    
    techs = [
        ("Flutter", "flutter", "Mobile App UI"),
        ("Supabase", "supabase", "Auth & Database"),
        ("Firebase", "firebase", "Crashlytics / Push"),
        ("NodeJS", "nodejs", "Backend Services"),
        ("OpenAI", "openai", "AI Text generation"),
        ("TomTom", "tomtom", "Route Calculation"),
        ("MapLibre GL", "maplibre", "Vector Maps"),
        ("Speech to Text", "voice", "Dictado por Voz"),
        ("Text to Speech", "audio", "Narración de Paradas")
    ]
    
    grid_cols = 3
    col_w = 180
    row_h = 160
    
    for idx, (tech_name, tech_key, tech_sub) in enumerate(techs):
        row = idx // grid_cols
        col = idx % grid_cols
        
        tx = 1230 + col * col_w + 10
        ty = 320 + row * row_h
        
        draw.rounded_rectangle([tx, ty, tx + 160, ty + 140], radius=14, fill=(15, 17, 24, 255), outline=PIL_BORDER, width=1)
        
        logo_path = os.path.join(TEMP_DIR, f"logo_{tech_key}.png")
        pasted_logo = False
        
        if os.path.exists(logo_path):
            try:
                logo_img = Image.open(logo_path).convert("RGBA")
                logo_img = logo_img.resize((48, 48), Image.Resampling.LANCZOS)
                img.paste(logo_img, (tx + 56, ty + 16), logo_img)
                pasted_logo = True
            except Exception as e:
                print(f"Error loading logo {tech_key}: {e}")
                
        if not pasted_logo:
            draw.ellipse([tx + 56, ty + 16, tx + 104, ty + 64], fill=(0, 122, 255, 30), outline=PIL_PRIMARY, width=2)
            if tech_key == "voice":
                draw.rounded_rectangle([tx + 74, ty + 28, tx + 86, ty + 46], radius=4, fill=PIL_PRIMARY)
                draw.line([tx + 80, ty + 46, tx + 80, ty + 54], fill=PIL_PRIMARY, width=2)
            elif tech_key == "audio":
                draw.polygon([(tx + 72, ty + 36), (tx + 82, ty + 36), (tx + 90, ty + 26), (tx + 90, ty + 54), (tx + 82, ty + 44), (tx + 72, ty + 44)], fill=PIL_PRIMARY)
            else:
                letter = tech_name[0]
                font_let = get_font("bold", 22)
                draw.text((tx + 70, ty + 26), letter, font=font_let, fill=PIL_TEXT_PRIMARY)
                
        font_tech_name = get_font("bold", 13)
        font_tech_sub = get_font("regular", 10)
        
        nw = font_tech_name.getlength(tech_name)
        draw.text((tx + 80 - nw // 2, ty + 78), tech_name, font=font_tech_name, fill=PIL_TEXT_PRIMARY)
        
        sw = font_tech_sub.getlength(tech_sub)
        draw.text((tx + 80 - sw // 2, ty + 104), tech_sub, font=font_tech_sub, fill=PIL_TEXT_SECONDARY)
        
    img.save(os.path.join(TEMP_DIR, "render_slide_2.png"))

# =========================================================================
# PPTX BUILDER (python-pptx)
# =========================================================================

def build_editable_pptx():
    """Build the widescreen editable PPTX presentation file overlaying text and images on backgrounds."""
    print("Building editable PowerPoint file...")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    # Define exact RGB Colors for PowerPoint API
    PPTX_BG = RGBColor(9, 11, 16)
    PPTX_CARD = RGBColor(28, 28, 30)
    PPTX_BORDER = RGBColor(56, 56, 58)
    PPTX_PRIMARY = RGBColor(0, 122, 255)
    PPTX_ACCENT = RGBColor(175, 82, 222)
    PPTX_SUCCESS = RGBColor(16, 185, 129)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 0: Portada
    # ----------------------------------------------------
    slide0 = prs.slides.add_slide(blank_layout)
    bg0_path = os.path.join(TEMP_DIR, "bg_cover.png")
    slide0.shapes.add_picture(bg0_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    
    logo_path = os.path.join(ASSETS_DIR, "logo_light.png")
    if os.path.exists(logo_path):
        slide0.shapes.add_picture(logo_path, Inches(6.04), Inches(1.94), width=Inches(1.25), height=Inches(1.25))
        
    tx_box = slide0.shapes.add_textbox(Inches(2.0), Inches(3.47), Inches(9.33), Inches(1.2))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "VibeTours."
    p.font.name = "Inter"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    tx_box2 = slide0.shapes.add_textbox(Inches(1.0), Inches(4.37), Inches(11.33), Inches(0.8))
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Planifica, crea y comparte experiencias impulsadas por Inteligencia Artificial."
    p2.font.name = "Inter"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(184, 184, 184)
    p2.alignment = PP_ALIGN.CENTER
    
    tx_box3 = slide0.shapes.add_textbox(Inches(4.5), Inches(4.93), Inches(4.33), Inches(0.4))
    tf3 = tx_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Apple Intelligence Standard"
    p3.font.name = "Inter"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = PPTX_PRIMARY
    p3.alignment = PP_ALIGN.CENTER
    
    # ----------------------------------------------------
    # SLIDE 1: El recorrido del usuario
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1_path = os.path.join(TEMP_DIR, "bg_roadmap.png")
    slide1.shapes.add_picture(bg1_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    
    title_box = slide1.shapes.add_textbox(Inches(0.55), Inches(0.42), Inches(10.0), Inches(0.8))
    tf_t = title_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = "El recorrido del usuario"
    p_t.font.name = "Inter"
    p_t.font.size = Pt(38)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 255, 255)
    
    sub_box = slide1.shapes.add_textbox(Inches(0.55), Inches(0.9), Inches(10.0), Inches(0.4))
    tf_s = sub_box.text_frame
    p_s = tf_s.paragraphs[0]
    p_s.text = "CÓMO FUNCIONA VIBETOURS"
    p_s.font.name = "Inter"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = PPTX_PRIMARY
    
    stations_data = [
        {
            "num": 1,
            "title": "Ingreso Seguro",
            "desc": "Inicia sesión con Google, Correo o accede directo mediante el Modo Demo sin registros.",
            "pos": "top",
            "mockup": "mockup_1_login.png",
            "sx": 180, "sy_m": 180, "sy_t": 620
        },
        {
            "num": 2,
            "title": "Crear Tour (Manual / IA)",
            "desc": "Planifica manual o usa el asistente de IA dictando tus preferencias con micrófono por voz.",
            "pos": "bottom",
            "mockup": "mockup_2_chat_ia.png",
            "sx": 480, "sy_m": 810, "sy_t": 200
        },
        {
            "num": 3,
            "title": "Publicar en Catálogo",
            "desc": "Los tours IA se guardan como borradores temporales y pasan a moderación antes de publicarse.",
            "pos": "top",
            "mockup": "mockup_3_publicar.png",
            "sx": 780, "sy_m": 180, "sy_t": 620
        },
        {
            "num": 4,
            "title": "Reportar Contenido",
            "desc": "Filtro de seguridad. Reporta spam, contenido inapropiado o rutas peligrosas para moderación.",
            "pos": "bottom",
            "mockup": "mockup_4_reportar.png",
            "sx": 1080, "sy_m": 810, "sy_t": 200
        },
        {
            "num": 5,
            "title": "Centro PQRS",
            "desc": "Soporte integrado en Supabase. Envía quejas o peticiones con respuesta en menos de 24h.",
            "pos": "top",
            "mockup": "mockup_5_pqrs.png",
            "sx": 1380, "sy_m": 180, "sy_t": 620
        },
        {
            "num": 6,
            "title": "Eliminación de Cuenta",
            "desc": "Garantía de privacidad. Controla tu cuenta, cierra sesión o elimina tu perfil permanentemente.",
            "pos": "bottom",
            "mockup": "mockup_6_perfil.png",
            "sx": 1680, "sy_m": 810, "sy_t": 200
        }
    ]
    
    mw_canvas = 253
    mh_canvas = 460
    tx_w = 300
    tx_h = 180
    
    for st in stations_data:
        sx_inch = Inches((st["sx"] / 1920.0) * 13.333)
        mock_w_inch = Inches((mw_canvas / 1920.0) * 13.333)
        mock_h_inch = Inches((mh_canvas / 1080.0) * 7.5)
        
        px_mockup = sx_inch - mock_w_inch / 2
        sy_px = st["sy_m"] - 40 if st["pos"] == "top" else st["sy_m"] - mh_canvas + 40
        py_mockup = Inches((sy_px / 1080.0) * 7.5)
        
        m_path = os.path.join(TEMP_DIR, st["mockup"])
        if os.path.exists(m_path):
            slide1.shapes.add_picture(m_path, px_mockup, py_mockup, width=mock_w_inch, height=mock_h_inch)
            
        card_w_inch = Inches((tx_w / 1920.0) * 13.333)
        card_h_inch = Inches((tx_h / 1080.0) * 7.5)
        card_x_inch = sx_inch - card_w_inch / 2
        card_y_inch = Inches((st["sy_t"] / 1080.0) * 7.5)
        
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x_inch, card_y_inch, card_w_inch, card_h_inch)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPTX_CARD
        shape.line.color.rgb = PPTX_BORDER
        shape.line.width = Pt(1.5)
        
        tb = slide1.shapes.add_textbox(card_x_inch + Inches(0.1), card_y_inch + Inches(0.05), card_w_inch - Inches(0.2), card_h_inch - Inches(0.1))
        tf_c = tb.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.1)
        tf_c.margin_right = Inches(0.1)
        tf_c.margin_top = Inches(0.1)
        
        p_ct = tf_c.paragraphs[0]
        p_ct.text = st["title"]
        p_ct.font.name = "Inter"
        p_ct.font.size = Pt(13)
        p_ct.font.bold = True
        p_ct.font.color.rgb = RGBColor(255, 255, 255)
        p_ct.space_after = Pt(6)
        
        p_cd = tf_c.add_paragraph()
        p_cd.text = st["desc"]
        p_cd.font.name = "Inter"
        p_cd.font.size = Pt(9.5)
        p_cd.font.color.rgb = RGBColor(184, 184, 184)
        
    # ----------------------------------------------------
    # SLIDE 2: Estado actual del sistema
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    bg2_path = os.path.join(TEMP_DIR, "bg_dashboard.png")
    slide2.shapes.add_picture(bg2_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    
    title2_box = slide2.shapes.add_textbox(Inches(0.55), Inches(0.42), Inches(10.0), Inches(0.8))
    tf_t2 = title2_box.text_frame
    p_t2 = tf_t2.paragraphs[0]
    p_t2.text = "Estado actual del sistema"
    p_t2.font.name = "Inter"
    p_t2.font.size = Pt(38)
    p_t2.font.bold = True
    p_t2.font.color.rgb = RGBColor(255, 255, 255)
    
    sub2_box = slide2.shapes.add_textbox(Inches(0.55), Inches(0.9), Inches(10.0), Inches(0.4))
    tf_s2 = sub2_box.text_frame
    p_s2 = tf_s2.paragraphs[0]
    p_s2.text = "ESTRUCTURA DE LÍMITES Y CAPA DE INTEGRACIÓN TÉCNICA"
    p_s2.font.name = "Inter"
    p_s2.font.size = Pt(16)
    p_s2.font.bold = True
    p_s2.font.color.rgb = PPTX_PRIMARY
    
    cards_data = [
        {
            "title": "Límite Demo IA",
            "desc": "Los usuarios invitados sin cuenta tienen un límite estricto de 2 generaciones de tours con IA. Al agotar el cupo de demostración local, el sistema solicita registrarse o iniciar sesión para continuar usando el planificador inteligente.",
            "x": 80, "y": 200, "w": 520, "h": 340,
            "accent": PPTX_ACCENT
        },
        {
            "title": "Dependencias de Mapas",
            "desc": "La cartografía y geocodificación se apoya en proveedores de fuentes abiertas: OpenStreetMap (OSM) para mapas, Photon y Nominatim para búsquedas de lugares, y Overpass API. Posee un mecanismo de contingencia para fallback automático.",
            "x": 80, "y": 570, "w": 520, "h": 340,
            "accent": PPTX_PRIMARY
        },
        {
            "title": "Moderación de Tours",
            "desc": "Los tours creados mediante IA quedan guardados como borradores privados. El administrador debe revisarlos y aprobarlos desde su panel restringido antes de que aparezcan en el catálogo público general, previniendo spam.",
            "x": 640, "y": 200, "w": 520, "h": 340,
            "accent": PPTX_PRIMARY
        },
        {
            "title": "Acceso Administrador",
            "desc": "El panel de administración cuenta con seguridad granular. El acceso está restringido a un único ID de usuario o correo configurado estrictamente en las variables de entorno de la base de datos de Supabase, bloqueando intrusos.",
            "x": 640, "y": 570, "w": 520, "h": 340,
            "accent": PPTX_ACCENT
        },
        {
            "title": "Tecnologías & Servicios Integrados",
            "desc": "",
            "x": 1200, "y": 200, "w": 640, "h": 710,
            "accent": PPTX_SUCCESS
        }
    ]
    
    for c in cards_data:
        x_inch = Inches((c["x"] / 1920.0) * 13.333)
        y_inch = Inches((c["y"] / 1080.0) * 7.5)
        w_inch = Inches((c["w"] / 1920.0) * 13.333)
        h_inch = Inches((c["h"] / 1080.0) * 7.5)
        
        sh = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_inch, y_inch, w_inch, h_inch)
        sh.fill.solid()
        sh.fill.fore_color.rgb = PPTX_CARD
        sh.line.color.rgb = PPTX_BORDER
        sh.line.width = Pt(1.5)
        
        bx_inch = Inches((4.0 / 1920.0) * 13.333)
        by_inch = Inches((40.0 / 1080.0) * 7.5)
        bx_pos = x_inch + Inches(0.2)
        by_pos = y_inch + Inches(0.28)
        
        bar_sh = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx_pos, by_pos, bx_inch, by_inch)
        bar_sh.fill.solid()
        bar_sh.fill.fore_color.rgb = c["accent"]
        bar_sh.line.fill.background()
        
        title_box = slide2.shapes.add_textbox(x_inch + Inches(0.32), y_inch + Inches(0.22), w_inch - Inches(0.4), Inches(0.6))
        tf_ct = title_box.text_frame
        p_ct = tf_ct.paragraphs[0]
        p_ct.text = c["title"]
        p_ct.font.name = "Inter"
        p_ct.font.size = Pt(17)
        p_ct.font.bold = True
        p_ct.font.color.rgb = RGBColor(255, 255, 255)
        
        if c["desc"]:
            body_box = slide2.shapes.add_textbox(x_inch + Inches(0.2), y_inch + Inches(0.9), w_inch - Inches(0.4), h_inch - Inches(1.0))
            tf_cb = body_box.text_frame
            tf_cb.word_wrap = True
            p_cb = tf_cb.paragraphs[0]
            p_cb.text = c["desc"]
            p_cb.font.name = "Inter"
            p_cb.font.size = Pt(11.5)
            p_cb.font.color.rgb = RGBColor(184, 184, 184)
            p_cb.line_spacing = 1.2
            
    techs = [
        ("Flutter", "flutter", "Mobile App UI"),
        ("Supabase", "supabase", "Auth & Database"),
        ("Firebase", "firebase", "Crashlytics / Push"),
        ("NodeJS", "nodejs", "Backend Services"),
        ("OpenAI", "openai", "AI Text generation"),
        ("TomTom", "tomtom", "Route Calculation"),
        ("MapLibre GL", "maplibre", "Vector Maps"),
        ("Speech to Text", "voice", "Dictado por Voz"),
        ("Text to Speech", "audio", "Narración de Paradas")
    ]
    
    grid_cols = 3
    col_w = 180
    row_h = 160
    
    for idx, (tech_name, tech_key, tech_sub) in enumerate(techs):
        row = idx // grid_cols
        col = idx % grid_cols
        
        tx = 1230 + col * col_w + 10
        ty = 320 + row * row_h
        
        tx_inch = Inches((tx / 1920.0) * 13.333)
        ty_inch = Inches((ty / 1080.0) * 7.5)
        tw_inch = Inches((160.0 / 1920.0) * 13.333)
        th_inch = Inches((140.0 / 1080.0) * 7.5)
        
        cell_sh = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx_inch, ty_inch, tw_inch, th_inch)
        cell_sh.fill.solid()
        cell_sh.fill.fore_color.rgb = RGBColor(15, 17, 24)
        cell_sh.line.color.rgb = PPTX_BORDER
        cell_sh.line.width = Pt(1.0)
        
        logo_path = os.path.join(TEMP_DIR, f"logo_{tech_key}.png")
        if os.path.exists(logo_path):
            lw_inch = Inches((48.0 / 1920.0) * 13.333)
            lh_inch = Inches((48.0 / 1080.0) * 7.5)
            lx_inch = tx_inch + tw_inch / 2 - lw_inch / 2
            ly_pos = ty_inch + Inches((16.0 / 1080.0) * 7.5)
            slide2.shapes.add_picture(logo_path, lx_inch, ly_pos, width=lw_inch, height=lh_inch)
            
        label_box = slide2.shapes.add_textbox(tx_inch, ty_inch + th_inch - Inches(0.65), tw_inch, Inches(0.6))
        tf_lbl = label_box.text_frame
        tf_lbl.word_wrap = True
        
        p_ln = tf_lbl.paragraphs[0]
        p_ln.text = tech_name
        p_ln.font.name = "Inter"
        p_ln.font.size = Pt(9.5)
        p_ln.font.bold = True
        p_ln.font.color.rgb = RGBColor(255, 255, 255)
        p_ln.alignment = PP_ALIGN.CENTER
        
        p_ls = tf_lbl.add_paragraph()
        p_ls.text = tech_sub
        p_ls.font.name = "Inter"
        p_ls.font.size = Pt(7.5)
        p_ls.font.color.rgb = RGBColor(184, 184, 184)
        p_ls.alignment = PP_ALIGN.CENTER
        
    path = OUTPUT_PPTX
    counter = 1
    base, ext = os.path.splitext(OUTPUT_PPTX)
    while True:
        try:
            prs.save(path)
            print(f"Presentation saved successfully: {path}")
            break
        except PermissionError:
            path = f"{base}_Updated_{counter}{ext}"
            counter += 1
            if counter > 20:
                print("Error: Could not save presentation, too many files locked.")
                break

# =========================================================================
# PDF BUILDER (PILLOW)
# =========================================================================

def compile_pdf():
    print("Compiling final PDF document...")
    slide0_img = Image.open(os.path.join(TEMP_DIR, "render_slide_0.png")).convert("RGB")
    slide1_img = Image.open(os.path.join(TEMP_DIR, "render_slide_1.png")).convert("RGB")
    slide2_img = Image.open(os.path.join(TEMP_DIR, "render_slide_2.png")).convert("RGB")
    
    path = OUTPUT_PDF
    counter = 1
    base, ext = os.path.splitext(OUTPUT_PDF)
    while True:
        try:
            slide0_img.save(
                path,
                save_all=True,
                append_images=[slide1_img, slide2_img]
            )
            print(f"PDF saved successfully: {path}")
            break
        except PermissionError:
            path = f"{base}_Updated_{counter}{ext}"
            counter += 1
            if counter > 20:
                print("Error: Could not save PDF, too many files locked.")
                break

# =========================================================================
# MAIN EXECUTION ROUTINE
# =========================================================================

if __name__ == "__main__":
    print("====================================================")
    print("VibeTours Presentation Builder started...")
    print("====================================================")
    
    process_all_screenshots()
    
    print("Generating background assets...")
    generate_cover_background()
    generate_roadmap_background()
    generate_dashboard_background()
    
    print("Rendering slides for PDF...")
    render_slide_0_image()
    render_slide_1_image()
    render_slide_2_image()
    
    build_editable_pptx()
    compile_pdf()
    
    print("====================================================")
    print("Success: Presentation resources successfully built.")
    print("====================================================")
